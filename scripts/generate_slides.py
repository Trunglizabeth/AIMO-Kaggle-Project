"""
Generate a simple PowerPoint from `docs/slide_notes.md` and images in `outputs/`.
Requires: python-pptx

Usage:
    python scripts/generate_slides.py --notes docs/slide_notes.md --images outputs --out outputs/sprint4_slides.pptx
"""
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
import os


def md_to_slides(md_path, prs):
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = [l.rstrip() for l in f]

    slide = None
    for line in lines:
        if not line:
            continue
        if line.startswith('# '):
            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = line[2:]
            subtitle.text = ''
        elif line.startswith('## '):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = line[3:]
        elif line.startswith('- '):
            if slide is None:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = ''
            body = slide.shapes.placeholders[1].text_frame
            p = body.add_paragraph()
            p.text = line[2:]
            p.level = 0
            p.font.size = Pt(18)
        else:
            # plain paragraph
            if slide is None:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = ''
            body = slide.shapes.placeholders[1].text_frame
            p = body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)


def add_images(prs, images_dir):
    imgs = [os.path.join(images_dir, f) for f in os.listdir(images_dir) if f.lower().endswith(('.png','.jpg','.jpeg'))]
    imgs.sort()
    for img in imgs:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        pic = slide.shapes.add_picture(img, left, top, width=Inches(8))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--notes', default='docs/slide_notes.md')
    parser.add_argument('--images', default='outputs')
    parser.add_argument('--out', default='outputs/sprint4_slides.pptx')
    args = parser.parse_args()

    prs = Presentation()
    md_to_slides(args.notes, prs)
    if os.path.isdir(args.images):
        add_images(prs, args.images)
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    prs.save(args.out)
    print('Saved', args.out)


if __name__ == '__main__':
    main()
